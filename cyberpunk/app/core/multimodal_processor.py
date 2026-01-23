"""
多模态文档处理器 - Multimodal Document Processor

功能：
- PDF文档解析和提取
- Word文档解析
- Excel表格提取
- PPT内容提取
- 图片OCR识别
- 表格结构化提取
- LaTeX公式解析
"""
from typing import Any, Dict, List, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass
from datetime import datetime
import io

from core.logging_config import get_logger

logger = get_logger(__name__)


@dataclass
class DocumentContent:
    """文档内容"""
    text: str
    images: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    formulas: List[str]
    metadata: Dict[str, Any]


class MultimodalDocumentProcessor:
    """
    多模态文档处理器
    
    支持的格式：
    - PDF: PyPDF2 + pdfplumber
    - Word: python-docx
    - Excel: openpyxl
    - PPT: python-pptx
    - 图片: pytesseract + paddleocr
    - 表格: camelot-py
    - 公式: sympy
    """
    
    def __init__(self):
        """初始化多模态处理器"""
        self.supported_formats = {
            "pdf": [".pdf"],
            "word": [".docx", ".doc"],
            "excel": [".xlsx", ".xls"],
            "ppt": [".pptx", ".ppt"],
            "image": [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif", ".webp"]
        }
        
        # 检查依赖库
        self._check_dependencies()
        
        logger.info("多模态文档处理器初始化完成")
    
    def _check_dependencies(self):
        """检查依赖库是否可用"""
        self.available_processors = {}
        
        # PDF处理
        try:
            import PyPDF2
            import pdfplumber
            self.available_processors["pdf"] = True
            logger.info("PDF处理器可用")
        except ImportError:
            self.available_processors["pdf"] = False
            logger.warning("PDF处理器不可用，请安装: pip install PyPDF2 pdfplumber")
        
        # Word处理
        try:
            import docx
            self.available_processors["word"] = True
            logger.info("Word处理器可用")
        except ImportError:
            self.available_processors["word"] = False
            logger.warning("Word处理器不可用，请安装: pip install python-docx")
        
        # Excel处理
        try:
            import openpyxl
            self.available_processors["excel"] = True
            logger.info("Excel处理器可用")
        except ImportError:
            self.available_processors["excel"] = False
            logger.warning("Excel处理器不可用，请安装: pip install openpyxl")
        
        # PPT处理
        try:
            import pptx
            self.available_processors["ppt"] = True
            logger.info("PPT处理器可用")
        except ImportError:
            self.available_processors["ppt"] = False
            logger.warning("PPT处理器不可用，请安装: pip install python-pptx")
        
        # OCR处理
        try:
            import pytesseract
            self.available_processors["ocr"] = True
            logger.info("OCR处理器可用")
        except ImportError:
            self.available_processors["ocr"] = False
            logger.warning("OCR处理器不可用，请安装: pip install pytesseract")
        
        # 表格提取
        try:
            import camelot
            self.available_processors["table"] = True
            logger.info("表格提取器可用")
        except ImportError:
            self.available_processors["table"] = False
            logger.warning("表格提取器不可用，请安装: pip install camelot-py[cv]")
    
    async def process_document(self, file_path: str) -> DocumentContent:
        """
        处理文档
        
        Args:
            file_path: 文档路径
        
        Returns:
            文档内容
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        # 根据文件扩展名选择处理器
        suffix = path.suffix.lower()
        
        if suffix == ".pdf":
            return await self._process_pdf(path)
        elif suffix in [".docx", ".doc"]:
            return await self._process_word(path)
        elif suffix in [".xlsx", ".xls"]:
            return await self._process_excel(path)
        elif suffix in [".pptx", ".ppt"]:
            return await self._process_ppt(path)
        elif suffix in self.supported_formats["image"]:
            return await self._process_image(path)
        else:
            raise ValueError(f"不支持的文件格式: {suffix}")
    
    async def _process_pdf(self, path: Path) -> DocumentContent:
        """处理PDF文档"""
        if not self.available_processors.get("pdf"):
            raise RuntimeError("PDF处理器不可用")
        
        import PyPDF2
        import pdfplumber
        
        text_content = []
        images = []
        tables = []
        
        # 使用PyPDF2提取文本
        try:
            with open(path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
        except Exception as e:
            logger.error(f"PyPDF2提取失败: {e}")
        
        # 使用pdfplumber提取表格
        try:
            with pdfplumber.open(path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # 提取表格
                    page_tables = page.extract_tables()
                    for table_idx, table in enumerate(page_tables):
                        tables.append({
                            "page": page_num + 1,
                            "table_index": table_idx,
                            "data": table,
                            "rows": len(table),
                            "cols": len(table[0]) if table else 0
                        })
        except Exception as e:
            logger.error(f"pdfplumber提取失败: {e}")
        
        # 提取图片（如果有OCR处理器）
        if self.available_processors.get("ocr"):
            # TODO: 实现PDF图片提取和OCR
            pass
        
        return DocumentContent(
            text="\n\n".join(text_content),
            images=images,
            tables=tables,
            formulas=[],
            metadata={
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "format": "pdf",
                "pages": len(text_content)
            }
        )
    
    async def _process_word(self, path: Path) -> DocumentContent:
        """处理Word文档"""
        if not self.available_processors.get("word"):
            raise RuntimeError("Word处理器不可用")
        
        import docx
        
        text_content = []
        images = []
        tables = []
        
        try:
            doc = docx.Document(path)
            
            # 提取段落文本
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # 提取表格
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    "table_index": table_idx,
                    "data": table_data,
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0
                })
            
            # 提取图片信息
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    images.append({
                        "type": "embedded_image",
                        "target": rel.target_ref
                    })
        
        except Exception as e:
            logger.error(f"Word文档处理失败: {e}")
        
        return DocumentContent(
            text="\n\n".join(text_content),
            images=images,
            tables=tables,
            formulas=[],
            metadata={
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "format": "word",
                "paragraphs": len(text_content),
                "tables": len(tables)
            }
        )
    
    async def _process_excel(self, path: Path) -> DocumentContent:
        """处理Excel文档"""
        if not self.available_processors.get("excel"):
            raise RuntimeError("Excel处理器不可用")
        
        import openpyxl
        
        text_content = []
        tables = []
        
        try:
            workbook = openpyxl.load_workbook(path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # 提取表格数据
                table_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        table_data.append([str(cell) if cell is not None else "" for cell in row])
                
                if table_data:
                    tables.append({
                        "sheet_name": sheet_name,
                        "data": table_data,
                        "rows": len(table_data),
                        "cols": len(table_data[0]) if table_data else 0
                    })
                    
                    # 将表格转换为文本描述
                    text_content.append(f"工作表: {sheet_name}")
                    text_content.append(f"行数: {len(table_data)}, 列数: {len(table_data[0]) if table_data else 0}")
        
        except Exception as e:
            logger.error(f"Excel文档处理失败: {e}")
        
        return DocumentContent(
            text="\n\n".join(text_content),
            images=[],
            tables=tables,
            formulas=[],
            metadata={
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "format": "excel",
                "sheets": len(tables)
            }
        )

    async def _process_ppt(self, path: Path) -> DocumentContent:
        """处理PPT文档"""
        if not self.available_processors.get("ppt"):
            raise RuntimeError("PPT处理器不可用")

        from pptx import Presentation

        text_content = []
        images = []
        tables = []

        try:
            prs = Presentation(path)

            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []

                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # 提取表格
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)

                        tables.append({
                            "slide": slide_idx + 1,
                            "data": table_data,
                            "rows": len(table_data),
                            "cols": len(table_data[0]) if table_data else 0
                        })

                    # 提取图片
                    if shape.shape_type == 13:  # Picture
                        images.append({
                            "slide": slide_idx + 1,
                            "type": "picture"
                        })

                if slide_text:
                    text_content.append(f"幻灯片 {slide_idx + 1}:\n" + "\n".join(slide_text))

        except Exception as e:
            logger.error(f"PPT文档处理失败: {e}")

        return DocumentContent(
            text="\n\n".join(text_content),
            images=images,
            tables=tables,
            formulas=[],
            metadata={
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "format": "ppt",
                "slides": len(text_content)
            }
        )

    async def _process_image(self, path: Path) -> DocumentContent:
        """处理图片（OCR识别）"""
        if not self.available_processors.get("ocr"):
            raise RuntimeError("OCR处理器不可用")

        import pytesseract
        from PIL import Image

        text_content = []

        try:
            # 打开图片
            image = Image.open(path)

            # OCR识别
            text = pytesseract.image_to_string(image, lang='eng+chi_sim')
            if text.strip():
                text_content.append(text)

            # 尝试提取表格（如果图片包含表格）
            # TODO: 使用更高级的表格识别算法

        except Exception as e:
            logger.error(f"图片OCR处理失败: {e}")

        return DocumentContent(
            text="\n".join(text_content),
            images=[{"path": str(path), "type": "source"}],
            tables=[],
            formulas=[],
            metadata={
                "file_name": path.name,
                "file_size": path.stat().st_size,
                "format": "image"
            }
        )

    def extract_formulas(self, text: str) -> List[str]:
        """
        提取LaTeX公式

        Args:
            text: 包含公式的文本

        Returns:
            公式列表
        """
        formulas = []

        # 匹配LaTeX公式模式
        import re

        # 行内公式 $...$
        inline_pattern = r'\$([^\$]+)\$'
        inline_matches = re.findall(inline_pattern, text)
        formulas.extend(inline_matches)

        # 行间公式 $$...$$
        display_pattern = r'\$\$([^\$]+)\$\$'
        display_matches = re.findall(display_pattern, text)
        formulas.extend(display_matches)

        # LaTeX环境 \begin{equation}...\end{equation}
        env_pattern = r'\\begin\{equation\}(.*?)\\end\{equation\}'
        env_matches = re.findall(env_pattern, text, re.DOTALL)
        formulas.extend(env_matches)

        return formulas

    def extract_tables_from_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        使用Camelot从PDF中提取表格

        Args:
            pdf_path: PDF文件路径

        Returns:
            表格列表
        """
        if not self.available_processors.get("table"):
            logger.warning("表格提取器不可用")
            return []

        try:
            import camelot

            # 提取表格
            tables = camelot.read_pdf(pdf_path, pages='all', flavor='lattice')

            result = []
            for idx, table in enumerate(tables):
                result.append({
                    "table_index": idx,
                    "page": table.page,
                    "data": table.df.values.tolist(),
                    "rows": len(table.df),
                    "cols": len(table.df.columns),
                    "accuracy": table.accuracy
                })

            return result

        except Exception as e:
            logger.error(f"Camelot表格提取失败: {e}")
            return []

